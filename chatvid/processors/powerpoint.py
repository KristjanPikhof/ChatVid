"""
PowerPoint file processor.

Extracts text from PowerPoint presentations (.pptx) using python-pptx.
Includes slide content, speaker notes, and table data.
"""

from pathlib import Path
from .base import DocumentProcessor, ProcessorRegistry


@ProcessorRegistry.register
class PowerPointProcessor(DocumentProcessor):
    """Processor for PowerPoint presentation files"""

    name = "PowerPoint Presentation"
    extensions = [".pptx"]
    dependencies = ["pptx"]  # python-pptx package

    def extract_text(self, file_path: Path) -> str:
        """Extract text from PowerPoint file.

        Extracts slide titles, body text, speaker notes, and table content.

        Args:
            file_path: Path to the PowerPoint file

        Returns:
            Extracted text from all slides
        """
        try:
            from pptx import Presentation

            prs = Presentation(str(file_path))
            text_parts = []

            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_text_parts = []

                # Add slide number header
                slide_text_parts.append(f"## Slide {slide_num}")

                # Extract title if exists
                if hasattr(slide, 'shapes') and slide.shapes.title:
                    try:
                        title_text = slide.shapes.title.text
                        if title_text:
                            slide_text_parts.append(f"\n### {title_text}\n")
                    except Exception:
                        pass

                # Extract text from all shapes
                shape_texts = []
                for shape in slide.shapes:
                    # Skip title (already extracted)
                    if hasattr(slide.shapes, 'title') and shape == slide.shapes.title:
                        continue

                    # Extract text from text frames
                    if hasattr(shape, 'text') and shape.text:
                        shape_texts.append(shape.text)

                    # Extract text from tables
                    if hasattr(shape, 'table'):
                        try:
                            table_text = self._extract_table_text(shape.table)
                            if table_text:
                                shape_texts.append(table_text)
                        except Exception:
                            pass

                if shape_texts:
                    slide_text_parts.append("\n".join(shape_texts))

                # Extract speaker notes
                if slide.has_notes_slide:
                    try:
                        notes_frame = slide.notes_slide.notes_text_frame
                        if notes_frame and notes_frame.text:
                            notes_text = notes_frame.text.strip()
                            if notes_text:
                                slide_text_parts.append(f"\n[Speaker Notes: {notes_text}]")
                    except Exception:
                        pass

                # Join slide parts and add to main text
                if len(slide_text_parts) > 1:  # More than just slide header
                    text_parts.append("\n".join(slide_text_parts))
                    text_parts.append("")  # Empty line between slides

            return "\n".join(text_parts)

        except Exception as e:
            import sys
            print(f"Error reading PowerPoint {file_path.name}: {e}", file=sys.stderr)
            return ""

    def _extract_table_text(self, table) -> str:
        """Extract text from a PowerPoint table.

        Args:
            table: python-pptx Table object

        Returns:
            Formatted table text with rows
        """
        try:
            rows_text = []
            for row in table.rows:
                cells_text = []
                for cell in row.cells:
                    if cell.text:
                        cells_text.append(cell.text.strip())
                if cells_text:
                    rows_text.append(" | ".join(cells_text))

            if rows_text:
                return "\n[Table]\n" + "\n".join(rows_text)
            return ""
        except Exception:
            return ""

    def get_metadata(self, file_path: Path) -> dict:
        """Extract metadata from PowerPoint.

        Args:
            file_path: Path to the PowerPoint file

        Returns:
            Dictionary with metadata (slide count, has notes, etc.)
        """
        try:
            from pptx import Presentation

            prs = Presentation(str(file_path))
            metadata = {
                'slide_count': len(prs.slides)
            }

            # Check if any slides have notes
            has_notes = any(slide.has_notes_slide for slide in prs.slides)
            metadata['has_notes'] = has_notes

            # Extract core properties if available
            if prs.core_properties:
                if prs.core_properties.title:
                    metadata['title'] = prs.core_properties.title
                if prs.core_properties.author:
                    metadata['author'] = prs.core_properties.author

            return metadata

        except Exception:
            return {}
